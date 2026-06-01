"""Render a Markdown deck (`---`-separated slides) into a styled .pptx.

Usage:
    uvx --from python-pptx --from markdown-it-py \\
        python docs/marketing/_render_deck.py docs/marketing/deck_executive.md

Slides are separated by `---` lines. Each slide:
  * First `#` line  → slide title
  * `## …` lines    → section headings
  * `> Notes: …`    → speaker notes (paragraph mode)
  * Tables          → real PPTX tables with the brand palette
  * Code fences     → monospace text boxes
  * Bullets         → bulleted body
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# Brand palette ---------------------------------------------------------------
BRAND_DARK = RGBColor(0x0C, 0x23, 0x40)     # title navy
BRAND_ACCENT = RGBColor(0x2C, 0x6E, 0xCB)   # accent blue
BRAND_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
BRAND_MUTED = RGBColor(0x5C, 0x6B, 0x7A)
BRAND_RULE = RGBColor(0xD0, 0xD7, 0xDE)
BRAND_TABLE_HDR = BRAND_DARK
BRAND_TABLE_BG_ALT = RGBColor(0xF6, 0xF8, 0xFA)
BRAND_CODE_BG = RGBColor(0xF6, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide geometry --------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.55)
MARGIN_TOP = Inches(0.55)
BODY_TOP = Inches(1.55)
BODY_HEIGHT = Inches(5.5)
BODY_WIDTH = Inches(13.333 - 1.1)


# ---------------------------------------------------------------------------
# Markdown parsing (small, just enough for our deck shape).
# ---------------------------------------------------------------------------


def split_slides(md: str) -> list[str]:
    """Split on `---` lines (but NOT on `---` inside fenced code blocks)."""
    out: list[str] = []
    buf: list[str] = []
    in_code = False
    for raw in md.splitlines():
        line = raw.rstrip()
        if line.startswith("```"):
            in_code = not in_code
            buf.append(line)
            continue
        if not in_code and line.strip() == "---":
            out.append("\n".join(buf).strip())
            buf = []
            continue
        buf.append(line)
    if buf:
        out.append("\n".join(buf).strip())
    return [s for s in out if s]


def strip_html_comments(md: str) -> str:
    return re.sub(r"<!--.*?-->", "", md, flags=re.DOTALL)


def parse_slide(md: str) -> dict:
    """Pull a slide's elements out of its markdown chunk.

    Returns:
        {
            "title": str | None,
            "subtitle": str | None,
            "blocks": list[block],   # each block: {kind, payload}
            "notes": str | None,
        }
    where block kinds are: "para", "bullets", "table", "code", "heading2".
    """
    title: str | None = None
    subtitle: str | None = None
    notes: list[str] = []
    blocks: list[dict] = []

    lines = md.splitlines()
    i = 0
    in_code = False
    code_buf: list[str] = []

    def flush_code() -> None:
        nonlocal code_buf
        if code_buf:
            blocks.append({"kind": "code", "lines": code_buf[:]})
            code_buf = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                flush_code()
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_buf.append(line)
            i += 1
            continue

        if title is None and stripped.startswith("# "):
            title = stripped[2:].strip()
            i += 1
            # Look for an immediately following "##" line as subtitle
            while i < len(lines) and not lines[i].strip():
                i += 1
            if i < len(lines) and lines[i].strip().startswith("## "):
                subtitle = lines[i].strip()[3:].strip()
                i += 1
            continue

        if stripped.startswith("> Notes:") or stripped.startswith(">Notes:"):
            note = stripped.split("Notes:", 1)[1].strip()
            if note:
                notes.append(note)
            # Consume following continuation `> …` lines too.
            i += 1
            while i < len(lines) and lines[i].lstrip().startswith(">"):
                cont = lines[i].lstrip()[1:].strip()
                if cont:
                    notes.append(cont)
                i += 1
            continue

        if stripped.startswith("## "):
            blocks.append({"kind": "heading2", "text": stripped[3:].strip()})
            i += 1
            continue

        if stripped.startswith("|") and "|" in stripped[1:]:
            # Start of a table — collect until a non-pipe line.
            tbl: list[str] = [stripped]
            i += 1
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                tbl.append(lines[i].strip())
                i += 1
            blocks.append({"kind": "table", "rows": _parse_table(tbl)})
            continue

        if stripped.startswith("- ") or stripped.startswith("* "):
            bullets: list[str] = []
            while i < len(lines) and (
                lines[i].lstrip().startswith("- ") or lines[i].lstrip().startswith("* ")
            ):
                stripped_b = lines[i].lstrip()
                bullets.append(stripped_b[2:].strip())
                i += 1
            blocks.append({"kind": "bullets", "items": bullets})
            continue

        if not stripped:
            i += 1
            continue

        # Plain paragraph — accumulate until blank line / section change.
        para: list[str] = [stripped]
        i += 1
        while i < len(lines):
            nxt = lines[i].strip()
            if not nxt:
                break
            if nxt.startswith(("- ", "* ", "## ", "|", "> ", "```", "# ")):
                break
            para.append(nxt)
            i += 1
        blocks.append({"kind": "para", "text": " ".join(para)})

    flush_code()
    return {
        "title": title,
        "subtitle": subtitle,
        "blocks": blocks,
        "notes": "\n\n".join(notes) if notes else None,
    }


def _parse_table(lines: list[str]) -> list[list[str]]:
    """Markdown table → list of rows. Drops the `---`-separator row."""
    rows: list[list[str]] = []
    for line in lines:
        cells = [c.strip() for c in line.strip("|").split("|")]
        # Drop the alignment separator row (`---`, `:--:`, etc.)
        if cells and all(re.match(r"^:?-{2,}:?$", c) for c in cells if c):
            continue
        rows.append(cells)
    return rows


# ---------------------------------------------------------------------------
# Rendering primitives.
# ---------------------------------------------------------------------------


def _set_run_basic(run, *, size=14, bold=False, color=BRAND_TEXT, mono=False) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Menlo" if mono else "Helvetica Neue"


def _render_inline(paragraph, text: str, *, size=14, color=BRAND_TEXT) -> None:
    """Tokenise simple **bold** + `code` + plain text into runs."""
    pattern = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`)")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos : m.start()]
            _set_run_basic(run, size=size, color=color)
        chunk = m.group(0)
        run = paragraph.add_run()
        if chunk.startswith("**"):
            run.text = chunk[2:-2]
            _set_run_basic(run, size=size, color=color, bold=True)
        else:  # `code`
            run.text = chunk[1:-1]
            _set_run_basic(run, size=size, color=BRAND_DARK, mono=True)
        pos = m.end()
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]
        _set_run_basic(run, size=size, color=color)


def _add_textbox(
    slide,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
) -> "object":
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    return tb


def _add_rect(slide, *, left, top, width, height, fill=BRAND_DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Slide layouts.
# ---------------------------------------------------------------------------


def render_title_slide(prs, slide_data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Full-bleed brand background
    _add_rect(slide, left=0, top=0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT, fill=BRAND_DARK)

    # Title block
    tb = _add_textbox(
        slide,
        left=Inches(0.7),
        top=Inches(2.6),
        width=Inches(12),
        height=Inches(2.4),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _render_inline(p, slide_data["title"] or "", size=54, color=WHITE)

    if slide_data.get("subtitle"):
        p2 = tb.text_frame.add_paragraph()
        p2.space_before = Pt(12)
        _render_inline(p2, slide_data["subtitle"], size=24, color=RGBColor(0xCD, 0xDB, 0xEC))

    # Body blocks become a tagline / brief paragraph.
    body_text = " ".join(
        b["text"] for b in slide_data["blocks"] if b["kind"] == "para"
    ).strip()
    if body_text:
        tb_body = _add_textbox(
            slide,
            left=Inches(0.7),
            top=Inches(5.4),
            width=Inches(12),
            height=Inches(1.6),
        )
        p3 = tb_body.text_frame.paragraphs[0]
        _render_inline(p3, body_text, size=16, color=RGBColor(0xE3, 0xEA, 0xF2))

    # Accent rule
    _add_rect(
        slide,
        left=Inches(0.7),
        top=Inches(2.45),
        width=Inches(1.2),
        height=Emu(34290),
        fill=BRAND_ACCENT,
    )

    if slide_data.get("notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def render_content_slide(prs, slide_data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide_data.get("title") or ""

    # Title bar
    tb = _add_textbox(
        slide,
        left=MARGIN_X,
        top=MARGIN_TOP,
        width=BODY_WIDTH,
        height=Inches(0.7),
    )
    p = tb.text_frame.paragraphs[0]
    _render_inline(p, title, size=32, color=BRAND_DARK)
    p.runs[0].font.bold = True

    # Accent rule under title
    _add_rect(
        slide,
        left=MARGIN_X,
        top=Inches(1.25),
        width=Inches(0.8),
        height=Emu(38100),
        fill=BRAND_ACCENT,
    )

    # Body: render blocks top→down inside the body region.
    cursor_y = BODY_TOP
    available_height = SLIDE_HEIGHT - cursor_y - Inches(0.4)
    body_blocks = slide_data["blocks"]
    if not body_blocks:
        return

    for block in body_blocks:
        if cursor_y > SLIDE_HEIGHT - Inches(0.5):
            break  # overflow — clip; pander's slide budget is the editor's job

        if block["kind"] == "heading2":
            h_tb = _add_textbox(
                slide,
                left=MARGIN_X,
                top=cursor_y,
                width=BODY_WIDTH,
                height=Inches(0.42),
            )
            hp = h_tb.text_frame.paragraphs[0]
            _render_inline(hp, block["text"], size=18, color=BRAND_DARK)
            hp.runs[0].font.bold = True
            cursor_y += Inches(0.55)

        elif block["kind"] == "para":
            text = block["text"]
            # Heuristic: estimate height from line count.
            est_lines = max(1, len(text) // 110 + 1)
            h = Inches(0.32 * est_lines + 0.15)
            ptb = _add_textbox(
                slide,
                left=MARGIN_X,
                top=cursor_y,
                width=BODY_WIDTH,
                height=h,
            )
            pp = ptb.text_frame.paragraphs[0]
            _render_inline(pp, text, size=14)
            cursor_y += h + Inches(0.05)

        elif block["kind"] == "bullets":
            items = block["items"]
            h = Inches(0.32 * len(items) + 0.2)
            btb = _add_textbox(
                slide,
                left=MARGIN_X,
                top=cursor_y,
                width=BODY_WIDTH,
                height=h,
            )
            for idx, item in enumerate(items):
                p = btb.text_frame.paragraphs[0] if idx == 0 else btb.text_frame.add_paragraph()
                p.space_after = Pt(4)
                p.level = 0
                # Bullet marker — python-pptx doesn't expose easy bullet
                # API, so we prefix with a glyph and indent.
                _render_inline(p, "•  " + item, size=14)
            cursor_y += h + Inches(0.1)

        elif block["kind"] == "table":
            rows = block["rows"]
            if not rows:
                continue
            n_cols = max(len(r) for r in rows)
            # Normalise row widths
            rows = [(r + [""] * (n_cols - len(r))) for r in rows]
            n_rows = len(rows)
            tbl_height = Inches(min(0.38 * n_rows + 0.1, 5.0))
            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                MARGIN_X, cursor_y, BODY_WIDTH, tbl_height,
            )
            tbl = tbl_shape.table
            for r_idx, row in enumerate(rows):
                for c_idx, cell_text in enumerate(row):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.margin_left = Emu(45720)
                    cell.margin_right = Emu(45720)
                    cell.margin_top = Emu(22860)
                    cell.margin_bottom = Emu(22860)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = BRAND_TABLE_HDR
                        fg = WHITE
                        bold = True
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = (
                            BRAND_TABLE_BG_ALT if r_idx % 2 == 0 else WHITE
                        )
                        fg = BRAND_TEXT
                        bold = False
                    p = cell.text_frame.paragraphs[0]
                    # Clear default empty run by setting text first
                    p.text = ""
                    _render_inline(p, cell_text, size=11, color=fg)
                    for run in p.runs:
                        run.font.bold = bold or run.font.bold
            cursor_y += tbl_height + Inches(0.15)

        elif block["kind"] == "code":
            code_lines = block["lines"]
            line_count = len(code_lines)
            h = Inches(min(0.27 * line_count + 0.2, 5.0))
            # Background panel
            _add_rect(
                slide,
                left=MARGIN_X,
                top=cursor_y,
                width=BODY_WIDTH,
                height=h,
                fill=BRAND_CODE_BG,
            )
            ctb = _add_textbox(
                slide,
                left=MARGIN_X + Inches(0.15),
                top=cursor_y + Inches(0.1),
                width=BODY_WIDTH - Inches(0.3),
                height=h - Inches(0.2),
            )
            for idx, line in enumerate(code_lines):
                p = ctb.text_frame.paragraphs[0] if idx == 0 else ctb.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line if line else " "
                _set_run_basic(run, size=10, color=BRAND_DARK, mono=True)
            cursor_y += h + Inches(0.15)

    # Speaker notes
    if slide_data.get("notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def render_section_divider(prs, slide_data: dict) -> None:
    """Slim divider style for slides whose first block is just a heading.
    Currently unused — content slides use this layout naturally."""
    render_content_slide(prs, slide_data)


# ---------------------------------------------------------------------------
# Top-level entry point.
# ---------------------------------------------------------------------------


def build(md_path: Path, out_path: Path) -> None:
    md = strip_html_comments(md_path.read_text(encoding="utf-8"))
    slides_md = split_slides(md)
    if not slides_md:
        raise SystemExit(f"No slides found in {md_path}")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for idx, chunk in enumerate(slides_md):
        data = parse_slide(chunk)
        if idx == 0:
            render_title_slide(prs, data)
        else:
            render_content_slide(prs, data)

    prs.save(str(out_path))
    print(f"  wrote {out_path}  ({len(slides_md)} slides)")


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("usage: _render_deck.py <input.md>", file=sys.stderr)
        return 2
    md_path = Path(argv[1]).resolve()
    out_path = md_path.with_suffix(".pptx")
    build(md_path, out_path)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
